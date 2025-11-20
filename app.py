from flask import Flask, render_template, request, jsonify, session
import os
import sys
import uuid
import json
from datetime import datetime
from werkzeug.utils import secure_filename

# Optional imports with error handling
try:
    import pytesseract
except ImportError:
    print("Warning: pytesseract not available - OCR features disabled")
    pytesseract = None

try:
    import PyPDF2
except ImportError:
    print("Warning: PyPDF2 not available - PDF processing disabled")
    PyPDF2 = None

try:
    from docx import Document
except ImportError:
    print("Warning: python-docx not available - DOCX processing disabled")
    Document = None

try:
    from PIL import Image
except ImportError:
    print("Warning: Pillow not available - image processing disabled")
    Image = None

try:
    from pptx import Presentation
except ImportError:
    print("Warning: python-pptx not available - PPT creation disabled")
    Presentation = None

import io
import base64
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Preformatted, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.colors import HexColor
import re
import json
from datetime import datetime

# Suppress GPT4All CUDA warnings
os.environ['CUDA_VISIBLE_DEVICES'] = ''

# AI Providers
from groq import Groq
import google.generativeai as genai
try:
    from gpt4all import GPT4All
except ImportError:
    GPT4All = None

# Configure Tesseract path (if needed)
try:
    import pytesseract
    # Uncomment if tesseract not in PATH:
    # pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
except ImportError:
    pass

app = Flask(__name__)
app.secret_key = 'your-secret-key-here'

# Configuration
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'txt', 'pdf', 'png', 'jpg', 'jpeg', 'gif', 'docx'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size

# Create upload and chat history directories
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs('chat_history', exist_ok=True)

# AI Configuration
AI_CONFIGS = {
    'groq': {
        'api_key': 'gsk_YOUR_GROQ_API_KEY_HERE',
        'models': ['llama-3.1-8b-instant', 'llama-3.1-70b-versatile', 'mixtral-8x7b-32768']
    },
    'gemini': {
        'api_key': 'YOUR_GEMINI_API_KEY_HERE',
        'models': ['gemini-2.0-flash', 'gemini-1.5-pro']
    },
    'local': {
        'models': ['Llama-3.2-1B-Instruct-Q4_0.gguf']
    }
}

class AIManager:
    def __init__(self):
        self.groq_client = None
        self.gemini_model = None
        self.local_model = None
        self.uploaded_file_content = None
        self.init_clients()
    
    def save_chat_message(self, chat_id, message, response, provider):
        """Save chat message to JSON file"""
        chat_file = f'chat_history/{chat_id}.json'
        
        # Load existing chat or create new
        if os.path.exists(chat_file):
            with open(chat_file, 'r', encoding='utf-8') as f:
                chat_data = json.load(f)
        else:
            chat_data = {
                'chat_id': chat_id,
                'created_at': datetime.now().isoformat(),
                'messages': []
            }
        
        # Add new message
        chat_data['messages'].append({
            'timestamp': datetime.now().isoformat(),
            'user_message': message,
            'ai_response': response,
            'provider': provider
        })
        
        # Save to file
        with open(chat_file, 'w', encoding='utf-8') as f:
            json.dump(chat_data, f, indent=2, ensure_ascii=False)
    
    def get_chat_history(self, chat_id):
        """Get chat history from JSON file"""
        chat_file = f'chat_history/{chat_id}.json'
        if os.path.exists(chat_file):
            with open(chat_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        return None
    
    def get_all_chats(self):
        """Get list of all chat sessions"""
        chats = []
        if os.path.exists('chat_history'):
            for filename in os.listdir('chat_history'):
                if filename.endswith('.json'):
                    chat_id = filename[:-5]  # Remove .json extension
                    chat_data = self.get_chat_history(chat_id)
                    if chat_data and chat_data.get('messages'):
                        first_message = chat_data['messages'][0]['user_message'][:50] + '...' if len(chat_data['messages'][0]['user_message']) > 50 else chat_data['messages'][0]['user_message']
                        chats.append({
                            'chat_id': chat_id,
                            'created_at': chat_data.get('created_at'),
                            'message_count': len(chat_data.get('messages', [])),
                            'first_message': first_message
                        })
        return sorted(chats, key=lambda x: x['created_at'], reverse=True)
    
    def init_clients(self):
        try:
            if AI_CONFIGS['groq']['api_key'] != 'gsk_YOUR_GROQ_API_KEY_HERE':
                self.groq_client = Groq(api_key=AI_CONFIGS['groq']['api_key'])
                print("SUCCESS: Groq client initialized")
        except Exception as e:
            print(f"ERROR: Groq init failed: {e}")
            pass
        
        try:
            if AI_CONFIGS['gemini']['api_key'] != 'YOUR_GEMINI_API_KEY_HERE':
                genai.configure(api_key=AI_CONFIGS['gemini']['api_key'])
                self.gemini_model = genai.GenerativeModel('gemini-2.0-flash')
                print("SUCCESS: Gemini client initialized")
        except Exception as e:
            print(f"ERROR: Gemini init failed: {e}")
            pass
        
        try:
            if GPT4All:
                # Use model from GPT4All default directory
                model_path = r"C:\Users\Hp\AppData\Local\nomic.ai\GPT4All"
                available_models = [
                    "Meta-Llama-3-8B-Instruct.Q4_0.gguf",
                    "Llama-3.2-1B-Instruct-Q4_0.gguf",
                    "tinyllama-1.1b-chat-v1.0.Q4_K_M.gguf",
                    "mistral-7b-instruct-v0.1.Q4_0.gguf"
                ]
                
                model_name = None
                for model in available_models:
                    if os.path.exists(os.path.join(model_path, model)):
                        model_name = model
                        break
                
                if model_name:
                    self.local_model = GPT4All(
                        model_name,
                        model_path=model_path,
                        verbose=False,
                        device='cpu',
                        n_threads=6,
                        n_ctx=1024
                    )
                    print(f"SUCCESS: GPT4All {model_name} initialized (optimized)")
                else:
                    print("No GPT4All models found in default directory")
        except Exception as e:
            print(f"ERROR: GPT4All init failed: {e}")
            self.local_model = None
    
    def get_response(self, message, provider='auto', model=None, chat_id=None):
        # Check for developer/name questions
        dev_keywords = ['who made', 'who created', 'who developed', 'who built', 'developer', 'creator', 'team', 'your name', 'what is your name', 'who are you']
        if any(keyword in message.lower() for keyword in dev_keywords):
            return "I am nova, an AI assistant developed by Team Avinash, Anand, Akshay, and Mudhura. We are a collaborative team working on advanced AI solutions."
        
        if provider == 'auto':
            if self.groq_client:
                provider = 'groq'
            elif self.gemini_model:
                provider = 'gemini'
            elif self.local_model:
                provider = 'local'
            else:
                return "No AI provider available. Please configure API keys or install GPT4All."
        
        # Skip conversation history for simple greetings to get clean responses
        conversation_history = []
        simple_messages = ['hi', 'hello', 'hey', 'good morning', 'good afternoon', 'good evening']
        if not any(simple in message.lower() for simple in simple_messages) and chat_id:
            chat_data = self.get_chat_history(chat_id)
            if chat_data and chat_data.get('messages'):
                # Get last 3 exchanges for context (reduced to avoid confusion)
                recent_messages = chat_data['messages'][-3:]
                for msg in recent_messages:
                    conversation_history.append({"role": "user", "content": msg['user_message']})
                    conversation_history.append({"role": "assistant", "content": msg['ai_response']})
        
        # Add file content to message if available
        if self.uploaded_file_content:
            message = f"Based on the uploaded file content:\n{self.uploaded_file_content}\n\nUser question: {message}"
        
        try:
            if provider == 'groq' and self.groq_client:
                model = model or 'llama-3.1-8b-instant'
                # Build messages with conversation history
                messages = conversation_history + [{"role": "user", "content": message}]
                response = self.groq_client.chat.completions.create(
                    messages=messages,
                    model=model,
                    max_tokens=1000
                )
                return response.choices[0].message.content
            
            elif provider == 'gemini' and self.gemini_model:
                # For simple messages, don't include context
                if any(simple in message.lower() for simple in ['hi', 'hello', 'hey']):
                    response = self.gemini_model.generate_content(message)
                else:
                    # For complex messages, include context
                    context_message = message
                    if conversation_history:
                        context = "\n".join([f"{msg['role'].title()}: {msg['content']}" for msg in conversation_history[-4:]])
                        context_message = f"Previous conversation:\n{context}\n\nCurrent question: {message}"
                    response = self.gemini_model.generate_content(context_message)
                return response.text
            
            elif provider == 'local' and self.local_model:
                # Better prompt formatting for complete responses
                formatted_prompt = f"### Human: {message}\n### Assistant:"
                
                response = self.local_model.generate(
                    formatted_prompt,
                    max_tokens=200,
                    temp=0.7,
                    top_k=40,
                    top_p=0.9,
                    repeat_penalty=1.1,
                    streaming=False
                )
                
                # Clean up the response
                clean_response = response.replace(formatted_prompt, '').strip()
                clean_response = clean_response.replace('Human:', '').replace('Assistant:', '').strip()
                
                return clean_response if clean_response else "I understand your message."
            
            else:
                return f"Provider '{provider}' not available."
        
        except Exception as e:
            return f"Error: {str(e)}"

ai_manager = AIManager()

def execute_code(code, language='python'):
    """Execute code with interactive terminal support"""
    try:
        import subprocess
        import tempfile
        import os
        
        if language == 'python':
            # Create temporary file
            with tempfile.NamedTemporaryFile(mode='w', suffix='.py', delete=False) as f:
                f.write(code)
                temp_file = f.name
            
            try:
                # Execute with interactive input support
                result = subprocess.run(
                    ['python', '-u', temp_file],
                    input='5\n7\n',  # Provide sample inputs
                    capture_output=True,
                    text=True,
                    timeout=15
                )
                
                if result.returncode == 0:
                    return {'success': True, 'output': result.stdout}
                else:
                    return {'success': False, 'error': result.stderr}
                    
            finally:
                os.unlink(temp_file)
                
        elif language == 'javascript' or language == 'js':
            with tempfile.NamedTemporaryFile(mode='w', suffix='.js', delete=False) as f:
                f.write(code)
                temp_file = f.name
            
            try:
                result = subprocess.run(
                    ['node', temp_file],
                    capture_output=True,
                    text=True,
                    timeout=15
                )
                
                if result.returncode == 0:
                    return {'success': True, 'output': result.stdout}
                else:
                    return {'success': False, 'error': result.stderr}
                    
            finally:
                os.unlink(temp_file)
                
        elif language == 'java':
            # Extract class name from code
            import re
            class_match = re.search(r'public\s+class\s+(\w+)', code)
            if not class_match:
                return {'success': False, 'error': 'No public class found in Java code'}
            
            class_name = class_match.group(1)
            
            with tempfile.NamedTemporaryFile(mode='w', suffix=f'.java', delete=False) as f:
                f.write(code)
                temp_file = f.name
            
            try:
                # Compile
                compile_result = subprocess.run(
                    ['javac', temp_file],
                    capture_output=True,
                    text=True,
                    timeout=15
                )
                
                if compile_result.returncode != 0:
                    return {'success': False, 'error': compile_result.stderr}
                
                # Run
                class_file = temp_file.replace('.java', '.class')
                result = subprocess.run(
                    ['java', '-cp', os.path.dirname(temp_file), class_name],
                    input='5\n7\n',
                    capture_output=True,
                    text=True,
                    timeout=15
                )
                
                if result.returncode == 0:
                    return {'success': True, 'output': result.stdout}
                else:
                    return {'success': False, 'error': result.stderr}
                    
            finally:
                os.unlink(temp_file)
                if os.path.exists(temp_file.replace('.java', '.class')):
                    os.unlink(temp_file.replace('.java', '.class'))
                    
        else:
            return {'success': False, 'error': f'Language {language} not supported. Supported: python, javascript, java'}
            
    except subprocess.TimeoutExpired:
        return {'success': False, 'error': 'Code execution timed out (15s limit)'}
    except Exception as e:
        return {'success': False, 'error': str(e)}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_file(file_path, filename):
    try:
        ext = filename.rsplit('.', 1)[1].lower()
        
        if ext == 'txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        
        elif ext == 'pdf':
            text = ""
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
            return text if text.strip() else "No text found in PDF"
        
        elif ext == 'docx':
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text if text.strip() else "No text found in document"
        
        elif ext in ['png', 'jpg', 'jpeg', 'gif']:
            print(f"Processing image: {filename}")
            image = Image.open(file_path)
            
            # Convert to RGB if needed
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # Try multiple OCR configurations
            configs = ['--psm 6', '--psm 3', '--psm 1']
            
            for config in configs:
                try:
                    text = pytesseract.image_to_string(image, config=config)
                    if text.strip():
                        print(f"OCR successful with config: {config}")
                        return text.strip()
                except Exception as ocr_error:
                    print(f"OCR failed with config {config}: {ocr_error}")
                    continue
            
            return "No text found in image or OCR failed"
        
        return "Unsupported file type"
    
    except Exception as e:
        print(f"File processing error: {str(e)}")
        return f"Error extracting text: {str(e)}"

@app.route('/')
def index():
    if 'chat_id' not in session:
        session['chat_id'] = str(uuid.uuid4())
    return render_template('index.html')

@app.route('/chat', methods=['POST'])
def chat():
    data = request.json
    message = data.get('message', '')
    provider = data.get('provider', 'auto')
    model = data.get('model', None)
    
    if not message:
        return jsonify({'error': 'No message provided'})
    
    # Check for image generation request
    if any(word in message.lower() for word in ['create image', 'generate image', 'make image', 'draw']) or message.lower().startswith('image'):
        try:
            import requests
            import time
            from flask import send_file
            
            # Extract prompt
            prompt = message.lower().replace('create image', '').replace('generate image', '').replace('make image', '').replace('draw', '').replace('image', '').strip()
            if not prompt:
                prompt = "abstract art"
            
            url = f"https://image.pollinations.ai/prompt/{prompt.replace(' ', '%20')}"
            response = requests.get(url, timeout=30)
            
            if response.status_code == 200:
                buf = io.BytesIO(response.content)
                buf.seek(0)
                filename = f"ai_image_{int(time.time())}.jpg"
                return send_file(buf, download_name=filename, as_attachment=True, mimetype='image/jpeg')
            else:
                return jsonify({'error': 'Image generation failed'}), 500
        except Exception as e:
            print(f"Image generation failed: {e}")
            return jsonify({'error': f'Image generation error: {str(e)}'}), 500
    
    # Check for advanced document request with pages
    doc_request = parse_pages_and_format(message)
    if doc_request:
        topic = doc_request['topic'] or "General Topic"
        pages = int(doc_request.get('pages', 1))
        fmt = doc_request['format']
        
        try:
            content = generate_multi_page_content(topic, pages, fmt)
            
            if fmt == 'ppt':
                return create_multipage_ppt(content, topic, pages)
            elif fmt == 'docx':
                return create_multipage_docx(content, topic, pages)
            elif fmt == 'pdf':
                return create_multipage_pdf(content, topic, pages)
        except Exception as e:
            print(f"Document generation error: {e}")
            return jsonify({'error': 'Document creation failed. Try again.'})
    else:
        # Get chat_id for conversation memory
        chat_id = session.get('chat_id')
        response = ai_manager.get_response(message, provider, model, chat_id)
        
        # Save to chat history
        ai_manager.save_chat_message(chat_id, message, response, provider)
    
    return jsonify({'response': response})

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file selected'})
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'})
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        extracted_text = extract_text_from_file(file_path, filename)
        
        # Clean up uploaded file
        os.remove(file_path)
        
        # Store file content in AI manager for context
        ai_manager.uploaded_file_content = extracted_text
        
        return jsonify({
            'success': True,
            'filename': filename,
            'message': f'File "{filename}" uploaded successfully. You can now ask questions about it.'
        })
    
    return jsonify({'error': 'Invalid file type'})

@app.route('/providers')
def get_providers():
    available_providers = {}
    
    if ai_manager.groq_client:
        available_providers['groq'] = AI_CONFIGS['groq']['models']
    if ai_manager.gemini_model:
        available_providers['gemini'] = AI_CONFIGS['gemini']['models']
    if ai_manager.local_model:
        available_providers['local'] = AI_CONFIGS['local']['models']
    
    return jsonify(available_providers)

@app.route('/chat-history/<chat_id>')
def get_chat_history(chat_id):
    chat_data = ai_manager.get_chat_history(chat_id)
    if chat_data:
        return jsonify(chat_data)
    return jsonify({'error': 'Chat not found'}), 404

@app.route('/all-chats')
def get_all_chats():
    chats = ai_manager.get_all_chats()
    return jsonify(chats)

@app.route('/new-chat', methods=['POST'])
def new_chat():
    new_chat_id = str(uuid.uuid4())
    session['chat_id'] = new_chat_id
    return jsonify({'chat_id': new_chat_id})

@app.route('/clear-file', methods=['POST'])
def clear_file():
    ai_manager.uploaded_file_content = None
    return jsonify({'success': True})

@app.route('/clear-all-history', methods=['POST'])
def clear_all_history():
    import shutil
    try:
        if os.path.exists('chat_history'):
            shutil.rmtree('chat_history')
            os.makedirs('chat_history', exist_ok=True)
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# Store interactive sessions
interactive_sessions = {}

# Store API keys
api_keys = {}
import secrets

def generate_api_key():
    return f"sk-{secrets.token_urlsafe(32)}"

# Add CORS headers for API access
@app.after_request
def after_request(response):
    response.headers.add('Access-Control-Allow-Origin', '*')
    response.headers.add('Access-Control-Allow-Headers', 'Content-Type,Authorization')
    response.headers.add('Access-Control-Allow-Methods', 'GET,PUT,POST,DELETE,OPTIONS')
    return response

@app.route('/run-code', methods=['POST'])
def run_code():
    data = request.json
    code = data.get('code', '')
    language = data.get('language', 'python')
    
    if not code:
        return jsonify({'error': 'No code provided'})
    
    result = execute_code(code, language)
    return jsonify(result)

@app.route('/run-interactive', methods=['POST'])
def run_interactive():
    data = request.json
    code = data.get('code', '')
    language = data.get('language', 'python')
    session_id = data.get('session_id', '')
    
    if not code:
        return jsonify({'error': 'No code provided'})
    
    try:
        import subprocess
        import tempfile
        import os
        import threading
        import queue
        
        # Create temporary file with appropriate extension
        extensions = {
            'python': '.py',
            'javascript': '.js', 
            'java': '.java',
            'cpp': '.cpp',
            'c': '.c',
            'html': '.html',
            'css': '.css'
        }
        ext = extensions.get(language, '.txt')
        
        # Ensure temp directory exists
        temp_dir = tempfile.gettempdir()
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)
        
        return jsonify({'error': 'Interactive mode not implemented'})
    except Exception as e:
        return jsonify({'error': str(e)})

@app.route('/list-api-keys', methods=['GET'])
def list_api_keys():
    """List all API keys"""
    keys_list = []
    for key, info in api_keys.items():
        keys_list.append({
            'key': key[:20] + '...',  # Show only first 20 chars
            'full_key': key,
            'model': info['model'],
            'name': info['name'],
            'created': info['created'],
            'usage_count': info['usage_count']
        })
    return jsonify({'api_keys': keys_list})

@app.route('/create-api-key', methods=['POST'])
def create_api_key():
    """Create API key for specific model"""
    data = request.json
    model = data.get('model', '')
    name = data.get('name', 'Unnamed Key')
    
    if not model:
        return jsonify({'error': 'Model is required'}), 400
    
    api_key = generate_api_key()
    api_keys[api_key] = {
        'model': model,
        'name': name,
        'created': datetime.now().isoformat(),
        'usage_count': 0
    }
    
    return jsonify({
        'api_key': api_key,
        'model': model,
        'name': name,
        'message': 'API key created successfully'
    })

@app.route('/generate-image', methods=['POST'])
def generate_image():
    """Direct image generation endpoint"""
    try:
        data = request.json
        prompt = data.get('prompt', 'cat')
        
        import requests
        import time
        from flask import send_file
        
        url = f"https://image.pollinations.ai/prompt/{prompt.replace(' ', '%20')}"
        response = requests.get(url, timeout=30)
        
        if response.status_code == 200:
            buf = io.BytesIO(response.content)
            buf.seek(0)
            filename = f"ai_image_{int(time.time())}.jpg"
            return send_file(buf, download_name=filename, as_attachment=True, mimetype='image/jpeg')
        else:
            return jsonify({'error': 'Image generation failed'}), 500
    except Exception as e:
        return jsonify({'error': f'Image creation failed: {str(e)}'}), 500

@app.route('/send-input', methods=['POST'])
def send_input():
    data = request.json
    session_id = data.get('session_id', '')
    user_input = data.get('input', '')
    
    if session_id not in interactive_sessions:
        return jsonify({'error': 'Session not found'})
    
    session = interactive_sessions[session_id]
    process = session['process']
    
    try:
        # Send input to process
        process.stdin.write(user_input + '\n')
        process.stdin.flush()
        
        # Read output
        import time
        time.sleep(0.2)  # Give process time to respond
        
        output = ''
        while True:
            try:
                import select
                ready, _, _ = select.select([process.stdout], [], [], 0.5)
                if ready:
                    line = process.stdout.readline()
                    if line:
                        output += line
                    else:
                        break
                else:
                    break
            except:
                break
        
        # Check if process is still running
        if process.poll() is None:
            return jsonify({'waiting_for_input': True, 'output': output})
        else:
            # Process finished
            stdout, stderr = process.communicate()
            output += stdout
            
            # Cleanup
            os.unlink(session['temp_file'])
            del interactive_sessions[session_id]
            
            if process.returncode == 0:
                return jsonify({'success': True, 'output': output})
            else:
                return jsonify({'success': False, 'output': output, 'error': stderr})
                
    except Exception as e:
        return jsonify({'error': str(e)})

# ---------------------------
# Advanced "N pages" Document Generator
# ---------------------------

def parse_pages_and_format(text):
    t = text.lower()
    
    # Must contain both a format keyword AND an action keyword
    format_keywords = ["ppt", "powerpoint", "slides", "presentation", "docx", "document", "pdf"]
    action_keywords = ["create", "make", "generate", "build", "prepare"]
    
    has_format = any(k in t for k in format_keywords)
    has_action = any(k in t for k in action_keywords)
    
    if not (has_format and has_action):
        return None
    
    fmt = None
    # Check PDF first since it's most specific
    if "pdf" in t:
        fmt = "pdf"
    elif any(k in t for k in ["ppt", "powerpoint", "slides", "presentation"]):
        fmt = "ppt"
    elif any(k in t for k in ["word", "docx"]):
        fmt = "docx"
    elif "document" in t:
        # Default to docx for generic "document" requests
        fmt = "docx"

    # Extract number from beginning of text (like "7 Report Document AI")
    pages = 5  # default
    numbers = re.findall(r'\b(\d+)\b', text)
    if numbers:
        pages = max(1, min(20, int(numbers[0])))

    # Clean topic by removing numbers and control words
    remove = ["make", "create", "generate", "prepare", "build", "ppt", "pdf", "docx", "word",
              "powerpoint", "presentation", "slides", "of", "on", "about", "for", "pages", "page", "report"]
    words = []
    for word in text.split():
        if not word.isdigit() and word.lower() not in remove:
            words.append(word)
    
    topic = " ".join(words).strip()
    topic = topic if topic else "General Topic"

    return {"format": fmt, "topic": topic, "pages": pages}

def request_ai_structured_content(topic, pages, mode):
    try:
        prompt = f"""Create a {pages}-slide presentation about {topic}. Format as follows:

Slide 1: [Title]
- [Bullet point 1]
- [Bullet point 2] 
- [Bullet point 3]

Slide 2: [Title]
- [Bullet point 1]
- [Bullet point 2]
- [Bullet point 3]

Continue for all {pages} slides. Make each slide informative and well-structured."""
        
        if ai_manager.groq_client:
            resp = ai_manager.groq_client.chat.completions.create(
                messages=[{"role":"user","content":prompt}],
                model="llama-3.1-8b-instant",
                max_tokens=1000
            )
            return parse_simple_ai_response(resp.choices[0].message.content, topic, pages)
        elif ai_manager.gemini_model:
            resp = ai_manager.gemini_model.generate_content(prompt)
            return parse_simple_ai_response(resp.text, topic, pages)
    except Exception as e:
        print("AI generation failed:", e)
    return None

def parse_simple_ai_response(text, topic, pages):
    """Parse AI response into structured format"""
    try:
        lines = text.split('\n')
        slides = []
        current_slide = None
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Check if it's a slide title
            if any(word in line.lower() for word in ['slide', 'title', '##', '#']) or line.endswith(':'):
                if current_slide:
                    slides.append(current_slide)
                current_slide = {'heading': line.replace('#', '').replace('Slide', '').replace(':', '').strip(), 'bullets': []}
            elif line.startswith('-') or line.startswith('•') or line.startswith('*'):
                if current_slide:
                    current_slide['bullets'].append(line[1:].strip())
        
        if current_slide:
            slides.append(current_slide)
        
        # Ensure we have enough slides
        while len(slides) < pages:
            slides.append({
                'heading': f'{topic} - Part {len(slides) + 1}',
                'bullets': [f'Key point about {topic}', f'Important aspect of {topic}', f'Benefits of {topic}']
            })
        
        return {
            'title': topic.title(),
            'subtitle': 'AI Generated Presentation',
            'summary': f'This presentation covers {pages} key aspects of {topic}.',
            'pages': slides[:pages]
        }
    except:
        return None

def fallback_generate_content(topic, pages):
    title = topic.title()
    subtitle = "AI Generated Document"
    
    # Create topic-specific content
    if 'leave letter' in topic.lower():
        # Create actual leave letter template
        content_map = {
            1: {"heading": "Leave Application Letter Template", "bullets": ["Date: [Insert Date]", "To: [Manager/HR Name]", "Subject: Application for Leave"]},
            2: {"heading": "Letter Body", "bullets": ["Dear [Manager Name],", "I am writing to request leave from [start date] to [end date]", "Reason: [Personal/Medical/Emergency]"]},
            3: {"heading": "Work Arrangements", "bullets": ["I have completed all urgent tasks", "[Colleague name] will handle my responsibilities", "I will be available via email if needed"]},
            4: {"heading": "Closing", "bullets": ["Thank you for your consideration", "Sincerely, [Your Name]", "Contact: [Phone/Email]"]}            
        }
    elif 'ai' in topic.lower() or 'artificial intelligence' in topic.lower():
        content_map = {
            1: {"heading": "Introduction to Artificial Intelligence", "bullets": ["Definition and core concepts", "Brief history of AI development", "Current applications in daily life"]},
            2: {"heading": "Types of AI Systems", "bullets": ["Narrow AI vs General AI", "Machine Learning fundamentals", "Deep Learning and Neural Networks"]},
            3: {"heading": "AI Applications", "bullets": ["Healthcare and medical diagnosis", "Autonomous vehicles and transportation", "Natural language processing"]},
            4: {"heading": "Benefits and Advantages", "bullets": ["Increased efficiency and automation", "Enhanced decision-making capabilities", "24/7 availability and consistency"]},
            5: {"heading": "Challenges and Limitations", "bullets": ["Ethical considerations and bias", "Job displacement concerns", "Data privacy and security issues"]}
        }
    else:
        content_map = {}
        for i in range(1, pages + 1):
            content_map[i] = {
                "heading": f"{title} - Section {i}",
                "bullets": [
                    f"Key concept {i}.1 about {topic}",
                    f"Important aspect {i}.2 of {topic}", 
                    f"Practical application {i}.3"
                ]
            }
    
    pages_list = []
    for i in range(1, min(pages + 1, 6)):  # Limit to 5 predefined slides
        if i in content_map:
            pages_list.append(content_map[i])
        else:
            pages_list.append({
                "heading": f"{title} - Additional Topic {i}",
                "bullets": [f"Point {i}.1 about {topic}", f"Point {i}.2 with details", f"Point {i}.3 summary"]
            })
    
    # Fill remaining pages if needed
    while len(pages_list) < pages:
        i = len(pages_list) + 1
        pages_list.append({
            "heading": f"{title} - Section {i}",
            "bullets": [f"Additional point {i}.1", f"Supporting detail {i}.2", f"Conclusion {i}.3"]
        })

    # Adjust subtitle based on content type
    if 'leave letter' in topic.lower():
        subtitle = "Leave Application Template"
    elif 'letter' in topic.lower():
        subtitle = "Letter Template"
    
    return {"title": title, "subtitle": subtitle, "pages": pages_list}

def generate_multi_page_content(topic, pages, mode):
    """
    Try AI structured content; fallback to programmatic content if needed.
    """
    ai_data = request_ai_structured_content(topic, pages, mode)
    if ai_data:
        # If AI returned fewer pages, expand using fallback logic
        if len(ai_data.get("pages", [])) < pages:
            fallback = fallback_generate_content(topic, pages)
            combined = ai_data.get("pages", []) + fallback["pages"][len(ai_data.get("pages", [])):]
            ai_data["pages"] = combined[:pages]
            ai_data.setdefault("title", topic.title())
            ai_data.setdefault("subtitle", "")
            ai_data.setdefault("summary", fallback["summary"])
        else:
            ai_data["pages"] = ai_data["pages"][:pages]
        return ai_data
    else:
        return fallback_generate_content(topic, pages)

def create_multipage_ppt(content, topic, pages):
    try:
        if not Presentation:
            return jsonify({'error': 'PowerPoint creation not available. Install python-pptx: pip install python-pptx'}), 500
        
        from pptx.util import Pt
        from flask import send_file

        prs = Presentation()
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = content.get("title", topic.title())
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = "AI Generated Presentation"

        # Content slides
        pages_data = content.get("pages", [])
        for i in range(pages):
            page = pages_data[i] if i < len(pages_data) else {
                "heading": f"Topic {i+1}: {topic}", 
                "bullets": [f"Key point {i+1}.1 about {topic}", f"Important aspect {i+1}.2", f"Additional detail {i+1}.3"]
            }
            
            # Create content slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            # Set title
            slide.shapes.title.text = page.get("heading", f"Slide {i+1}")
            
            # Add bullet points
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            
            bullets = page.get("bullets", [])
            for j, bullet_text in enumerate(bullets[:5]):
                if j == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = str(bullet_text).strip()
                p.level = 0
                p.font.size = Pt(20)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        filename = f"{topic.replace(' ', '_')}_slides.pptx"
        return send_file(buf, download_name=filename, as_attachment=True)
    except Exception as e:
        return jsonify({'error': f'PPT creation failed: {str(e)}'}), 500

def create_multipage_docx(content, topic, pages):
    try:
        from flask import send_file
        from docx.shared import Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        doc = Document()
        
        # Title page
        title_para = doc.add_heading(content.get("title", topic.title()), level=0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        subtitle_para = doc.add_paragraph("AI Generated Document")
        subtitle_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        date_para = doc.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
        date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_paragraph("")  # Spacing
        
        # Content sections
        pages_data = content.get("pages", [])
        for i in range(pages):
            page = pages_data[i] if i < len(pages_data) else {
                "heading": f"Section {i+1}: {topic}", 
                "bullets": [f"Key point {i+1}.1 about {topic}", f"Important aspect {i+1}.2", f"Additional information {i+1}.3"]
            }
            
            # Add section heading
            heading = doc.add_heading(page.get("heading", f"Section {i+1}"), level=1)
            
            # Add bullet points
            bullets = page.get("bullets", [])
            for bullet in bullets:
                p = doc.add_paragraph()
                p.style = 'List Bullet'
                p.add_run(str(bullet).strip())
            
            # Add spacing between sections
            if i < pages - 1:
                doc.add_paragraph("")

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        filename = f"{topic.replace(' ', '_')}_report.docx"
        return send_file(buf, download_name=filename, as_attachment=True)
    except Exception as e:
        return jsonify({'error': f'DOCX creation failed: {str(e)}'}), 500

def create_multipage_pdf(content, topic, pages):
    """
    Create PDF with exactly `pages` pages using ReportLab platypus.
    """
    try:
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from flask import send_file

        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=(8.5*inch, 11*inch), rightMargin=0.6*inch, leftMargin=0.6*inch, topMargin=0.6*inch, bottomMargin=0.6*inch)
        styles = getSampleStyleSheet()
        # tweak styles
        heading_style = styles['Heading1']
        body_style = styles['BodyText']
        note_style = ParagraphStyle('NoteStyle', parent=styles['Italic'], fontSize=9)

        elements = []
        elements.append(Paragraph(content.get("title", topic.title()), heading_style))
        if content.get("subtitle"):
            elements.append(Paragraph(content.get("subtitle"), styles['Normal']))
        elements.append(Spacer(1, 12))
        elements.append(Paragraph("<b>Summary</b>", styles['Heading2']))
        elements.append(Paragraph(content.get("summary", ""), body_style))
        elements.append(Spacer(1, 18))

        pages_data = content.get("pages", [])
        for i in range(pages):
            page = pages_data[i] if i < len(pages_data) else {"heading": f"Page {i+1}", "body": "", "bullets": [], "notes": ""}
            elements.append(Paragraph(page.get("heading", f"Page {i+1}"), styles['Heading2']))
            # Break body into paragraphs approx
            body_text = page.get("body", "")
            for para in re.split(r'\n\s*\n', body_text.strip()):
                if para.strip():
                    elements.append(Paragraph(para.strip(), body_style))
                    elements.append(Spacer(1, 6))
            if page.get("bullets"):
                for b in page.get("bullets", []):
                    elements.append(Paragraph(f"• {b}", body_style))
            if page.get("notes"):
                elements.append(Spacer(1, 6))
                elements.append(Paragraph("Notes: " + page.get("notes"), note_style))
            # ensure page break after each page except last
            if i < pages - 1:
                elements.append(PageBreak())

        doc.build(elements)
        buf.seek(0)
        filename = f"{topic.replace(' ', '_')}_{pages}p.pdf"
        return send_file(buf, download_name=filename, as_attachment=True)
    except Exception as e:
        return jsonify({'error': f'PDF creation failed: {str(e)}'}), 500

def create_image_response(message):
    """Generate image from prompt"""
    try:
        import requests
        import time
        from flask import send_file
        
        # Extract prompt from message
        prompt = message.lower().replace('create image', '').replace('generate image', '').replace('make image', '').replace('draw', '').strip()
        if not prompt:
            prompt = "abstract art"
        
        print(f"Generating image for prompt: {prompt}")
        
        url = f"https://image.pollinations.ai/prompt/{prompt.replace(' ', '%20')}"
        response = requests.get(url, timeout=30)
        
        print(f"API response status: {response.status_code}")
        
        if response.status_code == 200:
            buf = io.BytesIO(response.content)
            buf.seek(0)
            filename = f"ai_image_{int(time.time())}.jpg"
            return send_file(buf, download_name=filename, as_attachment=True, mimetype='image/jpeg')
        else:
            return jsonify({'error': 'Image generation failed'}), 500
    except Exception as e:
        print(f"Image generation error: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': f'Image creation failed: {str(e)}'}), 500

if __name__ == '__main__':
    print("="*50)
    print("ChatGPT Clone - Starting Server")
    print("="*50)
    print(f"Python version: {sys.version}")
    print(f"Flask app created successfully")
    print(f"Available AI providers:")
    if ai_manager.groq_client:
        print("  ✓ Groq API")
    if ai_manager.gemini_model:
        print("  ✓ Gemini API")
    if ai_manager.local_model:
        print("  ✓ Local GPT4All")
    if not (ai_manager.groq_client or ai_manager.gemini_model or ai_manager.local_model):
        print("  ⚠ No AI providers available - check API keys")
    
    print("\n✓ Web Interface: http://localhost:5000")
    print("✓ API Endpoints:")
    print("  - GET  /api/status")
    print("  - GET  /api/models")
    print("  - POST /api/chat/completions")
    print("  - POST /api/generate")
    print("  - GET  /debug-keys (debug)")
    print(f"\n✓ Active API Keys: {len(api_keys)}")
    for key, info in api_keys.items():
        print(f"  - {info['name']}: {key[:20]}...")
    print("\nPress Ctrl+C to stop the server")
    print("="*50)
    
    if len(api_keys) == 0:
        print("\n💡 Tip: Create API keys via the web interface to use the API endpoints")
        print("   Visit: http://localhost:5000 and click 'Create API Key'")
        print("="*50)
    
    try:
        app.run(debug=True, host='0.0.0.0', port=5000)
    except Exception as e:
        print(f"\n✗ Server failed to start: {e}")
        import traceback
        traceback.print_exc()
        input("Press Enter to exit...")